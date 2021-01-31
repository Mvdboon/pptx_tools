import os
import re
import pptx


def replaceStr(file, find, replace, fixyear=False):
    prs = pptx.Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for par in shape.text_frame.paragraphs:
                for run in par.runs:
                    for f in find:
                        run.text = run.text.replace(f, replace)

        for slidemaster in prs.slide_masters:
            for slide_layout in slidemaster.slide_layouts:
                for shape in slide_layout.shapes:
                    if not shape.has_text_frame:
                        continue
                    for par in shape.text_frame.paragraphs:
                        for run in par.runs:
                            for f in find:
                                run.text = run.text.replace(f, replace)
                for pholder in slide_layout.placeholders:
                    if not pholder.has_text_frame:
                        continue
                    for ph in pholder.text_frame.paragraphs:
                        for run in ph.runs:
                            for f in find:
                                run.text = run.text.replace(f, replace)

    outfolder = "fixedYear" if fixyear else "replacedStr"
    os.makedirs(f'output\\{outfolder}', exist_ok=True)
    prs.save('output\\{}\\{}'.format(outfolder, file.split("\\")[-1]))
