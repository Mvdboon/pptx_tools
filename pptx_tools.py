import os
import re
import sys
import glob
import json
import multiprocessing
from datetime import datetime
import pptx

import click
import comtypes.client
from PyPDF2.pdf import *


@click.group()
def main():
    """For explanation on the subcommands use them plus the help flag."""
    if not os.path.exists('input'):
        os.mkdir("input")
        click.echo("ERROR: Please place the files in the input folder.")


@main.command()
@click.option("-f", "--file", default=None, help="If you only want to convert one file, select it here.")
@click.option("-i", "--identifiers", default=["A.", "B.", "C."], show_default=True, help="List of strings that all have to be present for excluding the page.")
def rm_Pages(identifiers, file):
    """Remove pages based on list of strings present on PDF page."""
    if identifiers.__class__ == str:
        tmp = identifiers.split(',')
        identifiers = []
        for t in tmp:
            identifiers.append(t)
    if not file:
        for f in glob.glob("input/*.pdf"):
            removePagesBasedOnTextSub(f, identifiers)
    else:
        removePagesBasedOnTextSub(file, identifiers)


@main.command()
@click.option("-n", "--num-processes", default=4, show_default=True, help="number of parallel jobs.")
@click.option("-h", "--hidden", default=False, help="Include hidden slides in PDF.", type=bool)
@click.option("-f", "--file", default=None, help="If you only want to convert one file, select it here.")
def convert2pdf(num_processes, file, hidden):
    """Convert pptx to pdf. Default behaviour is to convert all pptx in the input folder. Output will be placed in output folder."""
    if not file:
        convertall(num_processes, hidden)
    else:
        pptx2pdfSub(file, hidden) if "pptx" in file else click.Abort(
            "Wrong file type")


@main.command()
@click.argument('find', required=True, nargs=1)
@click.argument('replace', required=True, nargs=1)
@click.option("-f", "--file", default=None, help="If you only want to process one file, select it here.")
# @click.option("-r", "--regex", is_flag=True, help="Use regex")
def string_replace(file, find, replace):
    "Replace FIND with REPLACE string in the file. Default behaviour is to process all pptx in the input folder. Output will be placed in output folder."
    if file:
        replaceStr(file, [find], replace, fixyear=True)
    else:
        for f in glob.glob("input/*.pptx"):
            replaceStr(f, [find], replace, fixyear=True)


@main.command()
@click.option("-f", "--file", default=None, help="If you only want to process one file, select it here.")
@click.option('--year', default=str(datetime.today().year), help=f"Change if you need a different year then {datetime.today().year}.", type=str)
def fix_year(file, year):
    "Replace old years in file(s) with YEAR. Default search string is '<year> Deloitte'"
    replacement = f"{year} Deloitte"
    if file:
        replaceStr(file, ["2016 Deloitte", "2017 Deloitte", "2018 Deloitte",
                          "2019 Deloitte", "2020 Deloitte", "2021 Deloitte"], year, fixyear=True)
    else:
        for f in glob.glob("input/*.pptx"):
            replaceStr(f, ["2016 Deloitte", "2017 Deloitte", "2018 Deloitte",
                           "2019 Deloitte", "2020 Deloitte", "2021 Deloitte"], year, fixyear=True)


def removePagesBasedOnTextSub(filename, identifiers):
    os.makedirs('output', exist_ok=True)
    os.makedirs('output/removedPages', exist_ok=True)

    with open(filename, "rb") as inputfile:
        fname = filename.split('\\')[-1]
        with open(f"output/removedPages/{fname}", "wb") as outputfile:
            output = PdfFileWriter()
            for page in PdfFileReader(inputfile).pages:
                text = page.extractText()
                if all(x in text for x in identifiers):
                    continue
                output.addPage(page)
            output.write(outputfile)


def convertall(num_processes, hidden):
    if not os.path.exists('input'):
        click.Abort("No input folder")
    p = multiprocessing.Pool(num_processes)
    for file in glob.glob("input/*.pptx"):
        p.apply_async(pptx2pdfSub, (file, hidden))
    p.close()
    p.join()


def pptx2pdfSub(inputFileName, hidden=False):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    hidden_setting = 1 if hidden else 0

    os.makedirs('output/converted2PDF', exist_ok=True)
    input = "{}\{}".format(os.path.dirname(
        os.path.realpath(__file__)), inputFileName)
    output = "{}\{}".format(os.path.dirname(os.path.realpath(
        __file__)), inputFileName.replace("pptx", "pdf").replace("input", "output\converted2PDF"))

    deck = powerpoint.Presentations.Open(input)
    deck.ExportAsFixedFormat(output, FixedFormatType=2,
                             PrintHiddenSlides=hidden_setting)

    deck.Close()
    powerpoint.Quit()


def replaceStr(file, find, replace, fixyear=False):
    prs = pptx.Presentation(file)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for par in shape.text_frame.paragraphs:
                for run in par.runs:
                    for f in find:
                        run.text = run.text.replace(f, replace)

        for slidemaster in prs.slide_masters:
            for slide_layout in slidemaster.slide_layouts:
                for shape in slide_layout.shapes:
                    if not shape.has_text_frame:
                        continue
                    for par in shape.text_frame.paragraphs:
                        for run in par.runs:
                            for f in find:
                                run.text = run.text.replace(f, replace)
                for pholder in slide_layout.placeholders:
                    if not pholder.has_text_frame:
                        continue
                    for ph in pholder.text_frame.paragraphs:
                        for run in ph.runs:
                            for f in find:
                                run.text = run.text.replace(f, replace)

    outfolder = "fixedYear" if fixyear else "replacedStr"
    os.makedirs(f'output\\{outfolder}', exist_ok=True)
    prs.save('output\\{}\\{}'.format(outfolder, file.split("\\")[-1]))


if __name__ == "__main__":
    multiprocessing.freeze_support()
    if getattr(sys, 'frozen', False):
        main()